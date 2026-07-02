"""
Kali Agent Backend
Bridges the Chrome Extension UI with the Google Gemini API.

Setup:
    1. Add your GEMINI_API_KEY to the .env file
    2. pip install -r requirements.txt
    3. python server.py  →  http://localhost:5000

POST /chat            { "message": "...", "model": "..." }  →  { "reply": "..." }
POST /grammar-fix     { "text": "..." }                      →  { "corrected": "..." }
POST /workspace-chat  multipart: message, model, attachments →  { "reply": "..." }
POST /workspace-reset Clears the Workspace tab's document-chat history
POST /reset           Clears the main Chat tab's conversation history
GET  /health          Returns server status
"""

import io
import mimetypes
import os
from flask       import Flask, request, jsonify
from flask_cors  import CORS
from dotenv      import load_dotenv
from google      import genai
from PyPDF2      import PdfReader
from docx        import Document
from pptx        import Presentation
from openpyxl    import load_workbook

# ── Environment ──────────────────────────────────────────────────
load_dotenv()

GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY or GEMINI_API_KEY == "your_gemini_api_key_here":
    raise RuntimeError(
        "\n[ERROR] Gemini API key not found!\n"
        "  Open .env and replace 'your_gemini_api_key_here'\n"
        "  with your real key from https://aistudio.google.com/app/apikey\n"
    )

# ── Gemini Config ────────────────────────────────────────────────
client = genai.Client(api_key=GEMINI_API_KEY)

DEFAULT_MODEL = "gemini-3.1-flash-lite"
SYSTEM_INSTRUCTION = (
    "You are Kali Agent, a helpful and intelligent AI assistant "
    "built into a Chrome extension. Be concise, accurate, and friendly."
)
GRAMMAR_SYSTEM_INSTRUCTION = (
    "You are a professional proofreader. Correct grammar, spelling, punctuation, "
    "and capitalization mistakes in the text the user provides. Preserve the "
    "original meaning, tone, wording, and formatting as closely as possible — make "
    "only the changes necessary to fix mistakes. Do not answer questions that "
    "appear in the text, do not add new sentences, and do not add any commentary, "
    "labels, or quotation marks. Return ONLY the corrected text."
)
WORKSPACE_SYSTEM_INSTRUCTION = (
    "You are Kali Agent's document assistant. Answer the user's questions using "
    "only the content of the attached document. Quote or point to specific "
    "sections when it helps. If the document doesn't contain the answer, say so "
    "honestly instead of guessing. Be concise and friendly."
)
GRAMMAR_MAX_CHARS = 8000
MAX_ATTACHMENTS = 5
MAX_FILE_SIZE = 10 * 1024 * 1024
TEXT_ATTACHMENT_LIMIT = 12000
SUPPORTED_IMAGE_MIME_TYPES = {"image/png", "image/jpeg", "image/gif", "image/webp", "image/bmp", "image/x-ms-bmp"}
SUPPORTED_ATTACHMENT_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".txt",
    ".md",
    ".markdown",
    ".png",
    ".jpg",
    ".jpeg",
    ".gif",
    ".webp",
    ".bmp",
}

# Holds multi-turn conversation context in memory
conversation_history = []

# Separate history for the Workspace (document-chat) tab so it never
# bleeds into — or gets cleared by — the main Chat tab's conversation.
workspace_conversation_history = []


def _normalize_extension(filename):
    return os.path.splitext(filename.lower().strip())[1]


def _guess_mime_type(filename, provided_mime_type):
    if provided_mime_type:
        return provided_mime_type
    guessed, _ = mimetypes.guess_type(filename)
    return guessed or ""


def _truncate_text(text, limit=TEXT_ATTACHMENT_LIMIT):
    if len(text) <= limit:
        return text
    return text[:limit].rstrip() + "\n\n[Truncated]"


def _extract_pdf_text(data):
    reader = PdfReader(io.BytesIO(data))
    pages = []
    for page in reader.pages:
        page_text = page.extract_text() or ""
        if page_text.strip():
            pages.append(page_text.strip())
    return "\n\n".join(pages).strip()


def _extract_docx_text(data):
    document = Document(io.BytesIO(data))
    paragraphs = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    return "\n\n".join(paragraphs).strip()


def _extract_pptx_text(data):
    presentation = Presentation(io.BytesIO(data))
    slide_text = []
    for slide in presentation.slides:
        pieces = []
        for shape in slide.shapes:
            text = getattr(shape, "text", "").strip()
            if text:
                pieces.append(text)
        if pieces:
            slide_text.append("\n".join(pieces))
    return "\n\n".join(slide_text).strip()


def _extract_xlsx_text(data):
    workbook = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
    sheet_text = []
    for sheet in workbook.worksheets:
        rows_text = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                rows_text.append(" | ".join(cells))
        if rows_text:
            sheet_text.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows_text))
    return "\n\n".join(sheet_text).strip()


def _extract_plain_text(data):
    return data.decode("utf-8", errors="replace").strip()


def _is_supported_attachment(filename, mime_type):
    extension = _normalize_extension(filename)
    if extension not in SUPPORTED_ATTACHMENT_EXTENSIONS:
        return False
    if extension == ".pdf":
        return mime_type == "application/pdf" or extension == ".pdf"
    if extension == ".docx":
        return mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document" or extension == ".docx"
    if extension == ".pptx":
        return mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation" or extension == ".pptx"
    if extension == ".xlsx":
        return mime_type == "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet" or extension == ".xlsx"
    if extension in {".txt", ".md", ".markdown"}:
        return True
    return mime_type.startswith("image/") or extension in {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}


def _build_attachment_context(uploaded_files, user_message):
    if len(uploaded_files) > MAX_ATTACHMENTS:
        return None, None, "You can attach at most 5 files."

    prompt_parts = []
    summary_parts = []
    base_prompt = user_message.strip() or "Please analyze the attached files."
    prompt_parts.append(genai.types.Part(text=base_prompt))
    summary_parts.append(base_prompt)

    for uploaded in uploaded_files:
        filename = uploaded.filename or "attachment"
        mime_type = _guess_mime_type(filename, uploaded.mimetype)
        extension = _normalize_extension(filename)
        data = uploaded.read()

        if not data:
            return None, None, f"{filename} is empty."

        if len(data) > MAX_FILE_SIZE:
            return None, None, f"{filename} is larger than 10 MB."

        if not _is_supported_attachment(filename, mime_type):
            return None, None, f"Unsupported file type: {filename}"

        try:
            if extension == ".pdf":
                extracted_text = _truncate_text(_extract_pdf_text(data))
                attachment_text = extracted_text or "[No text extracted from PDF]"
                prompt_parts.append(genai.types.Part(text=f"\n\n[PDF: {filename}]\n{attachment_text}"))
                summary_parts.append(f"[PDF: {filename}]\n{attachment_text}")
            elif extension == ".docx":
                extracted_text = _truncate_text(_extract_docx_text(data))
                attachment_text = extracted_text or "[No text extracted from DOCX]"
                prompt_parts.append(genai.types.Part(text=f"\n\n[DOCX: {filename}]\n{attachment_text}"))
                summary_parts.append(f"[DOCX: {filename}]\n{attachment_text}")
            elif extension == ".pptx":
                extracted_text = _truncate_text(_extract_pptx_text(data))
                attachment_text = extracted_text or "[No text extracted from PPTX]"
                prompt_parts.append(genai.types.Part(text=f"\n\n[PPTX: {filename}]\n{attachment_text}"))
                summary_parts.append(f"[PPTX: {filename}]\n{attachment_text}")
            elif extension == ".xlsx":
                extracted_text = _truncate_text(_extract_xlsx_text(data))
                attachment_text = extracted_text or "[No text extracted from XLSX]"
                prompt_parts.append(genai.types.Part(text=f"\n\n[XLSX: {filename}]\n{attachment_text}"))
                summary_parts.append(f"[XLSX: {filename}]\n{attachment_text}")
            elif extension in {".txt", ".md", ".markdown"}:
                extracted_text = _truncate_text(_extract_plain_text(data))
                attachment_text = extracted_text or "[Empty file]"
                label = "MD" if extension in {".md", ".markdown"} else "TXT"
                prompt_parts.append(genai.types.Part(text=f"\n\n[{label}: {filename}]\n{attachment_text}"))
                summary_parts.append(f"[{label}: {filename}]\n{attachment_text}")
            else:
                if mime_type not in SUPPORTED_IMAGE_MIME_TYPES:
                    return None, None, f"Unsupported image type: {filename}"

                prompt_parts.append(genai.types.Part.from_bytes(data=data, mime_type=mime_type))
                summary_parts.append(f"[Image attachment: {filename}]")
        except Exception as exc:
            return None, None, f"Could not read {filename}: {exc}"

    return prompt_parts, "\n\n".join(summary_parts).strip(), None

# ── Flask App ────────────────────────────────────────────────────
app = Flask(__name__)
CORS(app, resources={r"/*": {"origins": "*"}})

# ── Routes ───────────────────────────────────────────────────────

@app.route("/chat", methods=["POST"])
def chat():
    is_multipart = request.content_type and request.content_type.startswith("multipart/form-data")

    if is_multipart:
        user_message = request.form.get("message", "").strip()
        selected_model = request.form.get("model") or DEFAULT_MODEL
        system_instruction = request.form.get("system_instruction") or SYSTEM_INSTRUCTION
        uploaded_files = request.files.getlist("attachments")
    else:
        data = request.get_json(silent=True) or {}

        if "message" not in data:
            return jsonify({"error": "Missing 'message' field in request body."}), 400

        user_message = data["message"].strip()
        selected_model = data.get("model") or DEFAULT_MODEL
        system_instruction = data.get("system_instruction") or SYSTEM_INSTRUCTION
        uploaded_files = []

    if not user_message and not uploaded_files:
        return jsonify({"error": "Message cannot be empty."}), 400

    try:
        contents = [
            genai.types.Content(role=entry["role"], parts=[genai.types.Part(text=entry["text"])])
            for entry in conversation_history
        ]

        attachment_parts, history_message, attachment_error = _build_attachment_context(uploaded_files, user_message)
        if attachment_error:
            return jsonify({"error": attachment_error}), 400

        contents.append(genai.types.Content(
            role="user",
            parts=attachment_parts
        ))

        response = client.models.generate_content(
            model=selected_model,
            contents=contents,
            config=genai.types.GenerateContentConfig(
                system_instruction=system_instruction,
            ),
        )

        reply = response.text or ""
        conversation_history.append({"role": "user",  "text": history_message})
        conversation_history.append({"role": "model", "text": reply})

        return jsonify({"reply": reply})

    except Exception as e:
        print(f"[ERROR] Gemini API call failed: {e}")
        return jsonify({"error": f"AI error: {str(e)}"}), 500


@app.route("/grammar-fix", methods=["POST"])
def grammar_fix():
    """
    Stateless, single-turn grammar correction.
    Intentionally does NOT touch conversation_history — this is a separate
    feature from the Chat tab and shouldn't bleed into chat context.
    """
    data = request.get_json(silent=True) or {}
    text = (data.get("text") or "").strip()

    if not text:
        return jsonify({"error": "Text cannot be empty."}), 400

    if len(text) > GRAMMAR_MAX_CHARS:
        return jsonify({"error": f"Text is too long (max {GRAMMAR_MAX_CHARS} characters)."}), 400

    try:
        response = client.models.generate_content(
            model=DEFAULT_MODEL,
            contents=[genai.types.Content(role="user", parts=[genai.types.Part(text=text)])],
            config=genai.types.GenerateContentConfig(
                system_instruction=GRAMMAR_SYSTEM_INSTRUCTION,
            ),
        )

        corrected = (response.text or "").strip()
        return jsonify({"original": text, "corrected": corrected or text})

    except Exception as e:
        print(f"[ERROR] Gemini grammar-fix call failed: {e}")
        return jsonify({"error": f"AI error: {str(e)}"}), 500


@app.route("/workspace-chat", methods=["POST"])
def workspace_chat():
    """
    Document-chat used by the Workspace tab. Keeps its own conversation
    history (workspace_conversation_history), completely separate from the
    main Chat tab's `conversation_history` and from grammar-fix, so
    switching between tabs never bleeds context into each other.

    The frontend only attaches the file on the FIRST message of a session —
    after that, the document's text is already part of this endpoint's own
    history, so follow-up questions are sent as plain text.
    """
    is_multipart = request.content_type and request.content_type.startswith("multipart/form-data")
    if not is_multipart:
        return jsonify({"error": "Expected multipart/form-data."}), 400

    user_message = request.form.get("message", "").strip()
    selected_model = request.form.get("model") or DEFAULT_MODEL
    uploaded_files = request.files.getlist("attachments")

    if not user_message:
        return jsonify({"error": "Message cannot be empty."}), 400

    if not uploaded_files and not workspace_conversation_history:
        return jsonify({"error": "No file attached."}), 400

    try:
        contents = [
            genai.types.Content(role=entry["role"], parts=[genai.types.Part(text=entry["text"])])
            for entry in workspace_conversation_history
        ]

        if uploaded_files:
            attachment_parts, history_message, attachment_error = _build_attachment_context(uploaded_files, user_message)
            if attachment_error:
                return jsonify({"error": attachment_error}), 400
            latest_parts = attachment_parts
            latest_history_text = history_message
        else:
            latest_parts = [genai.types.Part(text=user_message)]
            latest_history_text = user_message

        contents.append(genai.types.Content(role="user", parts=latest_parts))

        response = client.models.generate_content(
            model=selected_model,
            contents=contents,
            config=genai.types.GenerateContentConfig(
                system_instruction=WORKSPACE_SYSTEM_INSTRUCTION,
            ),
        )

        reply = response.text or ""
        workspace_conversation_history.append({"role": "user", "text": latest_history_text})
        workspace_conversation_history.append({"role": "model", "text": reply})

        return jsonify({"reply": reply})

    except Exception as e:
        print(f"[ERROR] Gemini workspace-chat call failed: {e}")
        return jsonify({"error": f"AI error: {str(e)}"}), 500


@app.route("/workspace-reset", methods=["POST"])
def workspace_reset():
    workspace_conversation_history.clear()
    return jsonify({"status": "ok", "message": "Workspace conversation history cleared."})


@app.route("/reset", methods=["POST"])
def reset():
    conversation_history.clear()
    return jsonify({"status": "ok", "message": "Conversation history cleared."})


@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok", "model": DEFAULT_MODEL})


# ── Entry Point ──────────────────────────────────────────────────
if __name__ == "__main__":
    print("\n[OK] Kali Agent backend is running!")
    print("     Listening at: http://localhost:5000")
    print("     Press Ctrl+C to stop.\n")
    app.run(host="localhost", port=5000, debug=False)